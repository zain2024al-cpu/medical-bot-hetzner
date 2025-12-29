#!/usr/bin/env python3
"""
إنشاء عرض تقديمي للبوت الطبي باستخدام python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_medical_bot_presentation():
    """إنشاء عرض تقديمي شامل للبوت الطبي"""

    # إنشاء عرض تقديمي جديد
    prs = Presentation()

    # تعريف الألوان
    BLUE = RGBColor(0, 114, 198)
    GREEN = RGBColor(0, 153, 76)
    RED = RGBColor(220, 53, 69)
    YELLOW = RGBColor(255, 193, 7)
    WHITE = RGBColor(255, 255, 255)
    DARK_BLUE = RGBColor(0, 51, 102)

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "بوت التقارير الطبية"
    subtitle.text = "نظام شامل لإدارة التقارير الطبية\n\nالمطور: عمر\nالتاريخ: ديسمبر 2025"

    # تنسيق العنوان
    title_tf = title.text_frame
    title_tf.paragraphs[0].font.size = Pt(44)
    title_tf.paragraphs[0].font.color.rgb = BLUE
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 2: نظرة عامة
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'نظرة عامة على البوت'
    tf = body_shape.text_frame
    tf.text = 'ما هو البوت؟'

    p = tf.add_paragraph()
    p.text = '• بوت تليجرام متخصص في إدارة التقارير الطبية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• يساعد في تسجيل ومتابعة جميع الإجراءات الطبية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• واجهة سهلة الاستخدام باللغة العربية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• نظام أمان متقدم وحفظ تلقائي'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'الإحصائيات الأساسية:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• 98 مريض مسجل'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 20 مترجم معتمد'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• دعم جميع أنواع الإجراءات الطبية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 24/7 متوفر'
    p.level = 1

    # تنسيق العنوان
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = BLUE

    # Slide 3: الشاشة الرئيسية
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'الشاشة الرئيسية'
    tf = body_shape.text_frame
    tf.text = 'القائمة الرئيسية:'

    p = tf.add_paragraph()
    p.text = '📝 إضافة تقرير جديد'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '📝 تعديل التقارير'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '📊 عرض التقارير اليومية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '🕐 الجدول اليومي'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'الميزات الرئيسية:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• تنقل سهل بين الأقسام'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• أيقونات واضحة لكل وظيفة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تصميم متجاوب مع جميع الأجهزة'
    p.level = 1

    # Slide 4: نظام إضافة التقارير
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'نظام إضافة التقارير'
    tf = body_shape.text_frame
    tf.text = 'خطوات إضافة التقرير الجديد:'

    p = tf.add_paragraph()
    p.text = '1. اختيار نوع الإجراء'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '2. اختيار المترجم (إجباري)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '3. إدخال بيانات المريض'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '4. تحديد التاريخ والوقت'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '5. اختيار المستشفى والقسم'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '6. اختيار الطبيب'
    p.level = 1

    # Slide 5: المترجم الإجباري
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'الميزة الجديدة - المترجم الإجباري'
    tf = body_shape.text_frame
    tf.text = 'نظام المترجم الإجباري:'

    p = tf.add_paragraph()
    p.text = 'المميزات:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• إجباري في كل تقرير جديد'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• آخر خطوة قبل جمع البيانات'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• قائمة محددة من 20 مترجماً معتمداً'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تحقق تلقائي من صحة الاسم'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'قائمة المترجمين المعتمدين:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = 'مصطفى، واصل، نجم الدين، محمد علي، سعيد، مهدي، صبري، عزي، معتز، ادريس،'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'هاشم، ادم، زيد، عصام، عزالدين، حسن، زين العابدين، عبدالسلام، ياسر، يحيى'
    p.level = 1

    # Slide 6: نظام البحث الذكي
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'نظام البحث الذكي'
    tf = body_shape.text_frame
    tf.text = 'كيفية عمل البحث:'

    p = tf.add_paragraph()
    p.text = 'البحث عن المرضى:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '🔍 ابحث عن المريض واكتبه'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '↓ يكتب المستخدم أول حرف'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '↓ قائمة بالأسماء المطابقة تظهر فوراً'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '↓ اضغط على الاسم المطلوب'
    p.level = 1

    # Slide 7: أنواع الإجراءات
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'أنواع الإجراءات الطبية'
    tf = body_shape.text_frame
    tf.text = 'الأنواع المدعومة:'

    p = tf.add_paragraph()
    p.text = 'أ. استشارة جديدة:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• شكوى المريض'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تشخيص الطبيب'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• الفحوصات المطلوبة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'ب. متابعة في الرقود:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• سبب المتابعة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• حالة المريض'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تاريخ العودة'
    p.level = 1

    # Slide 8: التعديل والعرض
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'نظام التعديل والعرض'
    tf = body_shape.text_frame
    tf.text = 'تعديل التقارير:'

    p = tf.add_paragraph()
    p.text = 'طريقة البحث عن التقرير:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• البحث بالتاريخ أو اسم المريض'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'الحقول القابلة للتعديل:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• اسم المريض والتاريخ'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• المستشفى والقسم والطبيب'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تفاصيل الإجراء الطبي'
    p.level = 1

    # Slide 9: الجدول اليومي والإشعارات
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'الجدول اليومي والإشعارات'
    tf = body_shape.text_frame
    tf.text = 'الجدول اليومي:'

    p = tf.add_paragraph()
    p.text = 'الميزات:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• المرضى المقرر زيارتهم اليوم'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• المواعيد المجدولة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• المتابعات المطلوبة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'نظام الإشعارات:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• تذكيرات يومية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تنبيهات للمواعيد'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تقارير دورية'
    p.level = 1

    # Slide 10: الأمان والحماية
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'الأمان والحماية'
    tf = body_shape.text_frame
    tf.text = 'ميزات الأمان:'

    p = tf.add_paragraph()
    p.text = 'الحماية الأساسية:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• التحقق من هوية المستخدم'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تشفير جميع البيانات'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• نسخ احتياطي تلقائي'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'نظام النسخ الاحتياطي:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• حفظ تلقائي في Google Cloud'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• نسخ احتياطي يومي'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• استعادة فورية'
    p.level = 1

    # Slide 11: واجهة المستخدم
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'واجهة المستخدم'
    tf = body_shape.text_frame
    tf.text = 'تصميم الواجهة:'

    p = tf.add_paragraph()
    p.text = 'نظام الألوان:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• 🔵 أزرار زرقاء - العمليات الأساسية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 🟢 أزرار خضراء - التأكيد'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 🔴 أزرار حمراء - الإلغاء'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'أزرار التنقل:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• ⬅️ رجوع'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• ❌ إلغاء'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• ✅ تأكيد'
    p.level = 1

    # Slide 12: تجربة المستخدم
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'تجربة المستخدم'
    tf = body_shape.text_frame
    tf.text = 'سهولة الاستخدام:'

    p = tf.add_paragraph()
    p.text = 'المميزات الأساسية:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '1. واجهة بسيطة - أزرار واضحة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '2. إرشادات مباشرة - نصوص توضيحية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '3. بحث ذكي - اقتراحات فورية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '4. تدفق منطقي - خطوات متسلسلة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'أداء النظام:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• ⚡ استجابة سريعة (< 2 ثانية)'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 📱 متوافق مع جميع الأجهزة'
    p.level = 1

    # Slide 13: الدعم والمساعدة
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'الدعم والمساعدة'
    tf = body_shape.text_frame
    tf.text = 'نظام المساعدة:'

    p = tf.add_paragraph()
    p.text = 'الأوامر المتاحة:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '/help - للمساعدة والإرشادات'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '/report - للإبلاغ عن مشاكل فنية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '/status - لمعرفة حالة النظام'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'الدعم الفني:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• دعم 24/7 متوفر'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• حل المشاكل الفنية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تحديثات دورية'
    p.level = 1

    # Slide 14: الإحصائيات
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'الإحصائيات والمقاييس'
    tf = body_shape.text_frame
    tf.text = 'أداء النظام:'

    p = tf.add_paragraph()
    p.text = 'البيانات الحالية:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• 98 مريض مسجل في النظام'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 20 مترجم معتمد ومرخص'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• 8 أنواع إجراءات طبية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'معدلات الأداء:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• معدل استجابة: < 2 ثانية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• نسبة نجاح: 99.9%'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• رضا المستخدمين: 95%'
    p.level = 1

    # Slide 15: مستقبل النظام
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'مستقبل النظام'
    tf = body_shape.text_frame
    tf.text = 'التطويرات المستقبلية:'

    p = tf.add_paragraph()
    p.text = 'الميزات القادمة:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• تطبيق جوال مصاحب'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تقارير ذكية بالذكاء الاصطناعي'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• تكامل مع أنظمة المستشفيات'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'التحسينات المخططة:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• واجهة مستخدم محسنة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• أداء أسرع'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• ميزات أمان إضافية'
    p.level = 1

    # Slide 16: الخاتمة
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'الخاتمة'
    tf = body_shape.text_frame
    tf.text = 'الخلاصة:'

    p = tf.add_paragraph()
    p.text = 'بوت التقارير الطبية هو نظام شامل ومتطور يساعد العاملين في المجال الطبي على:'
    p.level = 0

    p = tf.add_paragraph()
    p.text = '• تنظيم وإدارة التقارير الطبية بكفاءة'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• متابعة حالة المرضى بدقة وموثوقية'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• توفير الوقت والجهد في التسجيل اليدوي'
    p.level = 1

    p = tf.add_paragraph()
    p.text = '• ضمان دقة البيانات وسلامتها'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'البوت مصمم ليكون رفيقاً موثوقاً للفرق الطبية، يجمع بين البساطة والقوة في خدمة الرعاية الصحية.'
    p.level = 0

    # Slide 17: شكر وتقدير
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "شكراً للانتباه!"
    subtitle.text = "للاستفسارات والأسئلة:\n\nتواصل مع المطور\nعمر - ديسمبر 2025\n\n🌟 شكراً لكم! 🌟"

    # تنسيق الشكر
    title_tf = title.text_frame
    title_tf.paragraphs[0].font.size = Pt(48)
    title_tf.paragraphs[0].font.color.rgb = GREEN
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # حفظ العرض التقديمي
    prs.save('Medical_Bot_Presentation.pptx')
    print("✅ تم إنشاء العرض التقديمي بنجاح!")

if __name__ == "__main__":
    create_medical_bot_presentation()
